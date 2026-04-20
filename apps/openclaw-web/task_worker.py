#!/usr/bin/env python3
import csv
import json
import os
import re
import sqlite3
import subprocess
import time
import urllib.error
import urllib.request
from datetime import datetime, timezone
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = Path(os.getenv("OPENCLAW_WEB_DATA_DIR", str(BASE_DIR / "data")))
STATIC_DIR = Path(os.getenv("OPENCLAW_WEB_STATIC_DIR", str(BASE_DIR / "static")))
IMAGES_DIR = STATIC_DIR / "images"
DB_PATH = str(DATA_DIR / "tasks.db")
APP_DB_PATH = str(DATA_DIR / "app.db")
ARTIFACTS_ROOT = os.getenv("OPENCLAW_WEB_ARTIFACTS_DIR", str(BASE_DIR / "artifacts" / "tasks"))
OPENCLAW_CMD = ["openclaw", "agent"]
HERMES_BRIDGE_RUN_URL = os.environ.get("HERMES_BRIDGE_RUN_URL", "http://127.0.0.1:8787/run")

IMAGE_KEYWORDS = [
    "生成图片", "生成图", "生成一幅", "画一幅", "画一张", "生成一张",
    "生成画", "给我画", "画一个", "生成一个", "画个", "画个图",
    "生成个图", "给我画一个", "帮我画", "帮我生成",
]

REPORT_KEYWORDS = [
    "报告", "研究", "分析", "调研", "方案", "规划", "行业", "市场",
    "趋势", "对比", "竞品", "总结", "复盘", "长文", "白皮书",
]

HERMES_KEYWORDS = [
    "服务器", "日志", "排查", "部署", "配置", "脚本", "代码",
    "调试", "抓取", "浏览器", "终端", "ssh", "docker", "nginx",
    "haproxy", "sqlite", "openclaw", "hermes",
]

ARTIFACT_HINTS = {
    "pdf": ["pdf", "打印版", "下载版", "报告"],
    "markdown": ["markdown", "md", "文档", "报告", "总结"],
    "csv": ["csv", "表格", "清单", "名单", "数据表"],
    "json": ["json", "结构化", "接口结果", "机器可读"],
    "xlsx": ["xlsx", "excel", "工作簿", "电子表格"],
    "docx": ["docx", "word", "可编辑文档"],
    "pptx": ["pptx", "ppt", "幻灯片", "演示文稿"],
    "image": ["图片", "海报", "配图", "示意图", "封面", "插图"],
}

ARTIFACT_FALLBACKS = {
    "xlsx": "csv",
    "pptx": "markdown",
}


def get_db(path):
    conn = sqlite3.connect(path, check_same_thread=False)
    conn.execute("PRAGMA busy_timeout = 5000")
    conn.row_factory = sqlite3.Row
    return conn


def markdown_to_pdf(md_path, task_id):
    """用 markdown + weasyprint 生成中文 PDF（内容完整）"""
    import markdown as md_lib
    import weasyprint

    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    pdf_path = os.path.join(artifact_dir, "report.pdf")
    try:
        with open(md_path, encoding="utf-8") as f:
            md_content = f.read()
        html_body = md_lib.markdown(md_content, extensions=["tables", "fenced_code"])
        html_full = (
            '<!DOCTYPE html>\n<html lang="zh-CN">\n<head>\n<meta charset="utf-8">\n'
            '<style>\n'
            'body{font-family:"Noto Serif CJK SC","AR PL UMing CN",serif;font-size:12pt;'
            'line-height:1.8;max-width:800px;margin:2cm auto;padding:0 1cm;color:#333;}\n'
            'h1{color:#1a1a2e;border-bottom:2px solid #5465ff;padding-bottom:8px;}\n'
            'h2{color:#16213e;margin-top:24px;border-left:4px solid #5465ff;padding-left:12px;}\n'
            'table{border-collapse:collapse;width:100%;margin:16px 0;}\n'
            'th,td{border:1px solid #ddd;padding:8px;text-align:center;}\n'
            'th{background:#f0f2ff;}\n'
            'code{background:#f5f5f5;padding:2px 6px;border-radius:4px;}\n'
            'blockquote{border-left:4px solid #5465ff;margin:16px 0;padding:8px 16px;color:#555;}\n'
            '</style>\n</head>\n<body>' + html_body + "</body>\n</html>"
        )
        weasyprint.HTML(string=html_full).write_pdf(pdf_path)
        if os.path.exists(pdf_path):
            print(f"[worker] PDF生成成功: {os.path.getsize(pdf_path)} bytes")
            return f"/api/tasks/{task_id}/download/report.pdf"
    except Exception as e:
        print(f"[worker] markdown转PDF失败: {e}")
    return None


def find_artifact_links(task_id):
    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    if not os.path.isdir(artifact_dir):
        return []
    links = []
    for fname in os.listdir(artifact_dir):
        fpath = os.path.join(artifact_dir, fname)
        if os.path.isfile(fpath):
            links.append(f"/api/tasks/{task_id}/download/{fname}")
    return links


def load_report_markdown(task_id):
    path = os.path.join(ARTIFACTS_ROOT, task_id, "report.md")
    if not os.path.exists(path):
        return None
    try:
        with open(path, encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"[worker] 读取 report.md 失败: {e}")
        return None


def detect_requested_artifacts(prompt: str, task: dict) -> list[str]:
    raw = str(prompt or "")
    raw_lower = raw.lower()
    wanted = []
    for artifact, keywords in ARTIFACT_HINTS.items():
        if any(kw in raw for kw in keywords) or any(kw in raw_lower for kw in [k.lower() for k in keywords]):
            wanted.append(artifact)

    if any(kw in raw for kw in REPORT_KEYWORDS):
        if "markdown" not in wanted:
            wanted.append("markdown")
        if "pdf" not in wanted:
            wanted.append("pdf")

    if task.get("attachment_kind") == "image" and "image" not in wanted and not wanted:
        wanted.append("markdown")

    return wanted or ["markdown"]


def normalize_requested_artifacts(wanted: list[str]) -> list[str]:
    normalized = []
    for item in wanted:
        if item not in normalized:
            normalized.append(item)

    for src, target in ARTIFACT_FALLBACKS.items():
        if src in normalized and target not in normalized:
            normalized.append(target)

    return normalized


def choose_executor(prompt: str, task: dict) -> str:
    raw = str(prompt or "")
    raw_lower = raw.lower()
    if any(kw in raw for kw in IMAGE_KEYWORDS):
        return "local_image"
    if task.get("attachment_kind") in ("file", "image"):
        return "hermes"
    if any(kw in raw for kw in HERMES_KEYWORDS) or any(kw in raw_lower for kw in [k.lower() for k in HERMES_KEYWORDS]):
        return "hermes"
    if any(kw in raw for kw in REPORT_KEYWORDS):
        return "hermes"
    return "local_openclaw"


def build_task_plan(task: dict) -> dict:
    prompt = task.get("prompt_text") or ""
    wanted_artifacts = normalize_requested_artifacts(detect_requested_artifacts(prompt, task))
    executor = choose_executor(prompt, task)
    is_report = any(kw in prompt for kw in REPORT_KEYWORDS)
    return {
        "executor": executor,
        "wanted_artifacts": wanted_artifacts,
        "is_report": is_report,
    }


def extract_fenced_block(text: str, language: str) -> str | None:
    pattern = rf"```{language}\s*(.*?)```"
    match = re.search(pattern, text or "", flags=re.DOTALL | re.IGNORECASE)
    if match:
        return match.group(1).strip()
    return None


def extract_markdown_table(text: str) -> str | None:
    lines = (text or "").splitlines()
    blocks = []
    current = []
    for line in lines:
        if "|" in line:
            current.append(line.rstrip())
        else:
            if len(current) >= 2:
                blocks.append(current[:])
            current = []
    if len(current) >= 2:
        blocks.append(current)

    for block in blocks:
        divider = block[1].replace("|", "").replace("-", "").replace(":", "").strip()
        if len(block) >= 2 and not divider:
            return "\n".join(block)
    return None


def markdown_table_to_rows(table_text: str) -> list[list[str]]:
    rows = []
    for idx, line in enumerate((table_text or "").splitlines()):
        parts = [cell.strip() for cell in line.strip().strip("|").split("|")]
        if not parts:
            continue
        if idx == 1:
            cleaned = [p.replace("-", "").replace(":", "").strip() for p in parts]
            if all(not cell for cell in cleaned):
                continue
        rows.append(parts)
    return rows


def save_csv_rows(path: str, rows: list[list[str]]):
    with open(path, "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerows(rows)


def maybe_create_xlsx_from_rows(path: str, rows: list[list[str]]) -> bool:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return False

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row in rows:
        ws.append(row)
    wb.save(path)
    return True


def maybe_create_docx(task_id: str, markdown_text: str) -> str | None:
    try:
        from docx import Document  # type: ignore
    except Exception:
        return None

    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    path = os.path.join(artifact_dir, "report.docx")
    doc = Document()
    in_code = False
    for raw_line in (markdown_text or "").splitlines():
        line = raw_line.rstrip()
        if line.startswith("```"):
            in_code = not in_code
            continue
        if not line:
            doc.add_paragraph("")
            continue
        if in_code:
            doc.add_paragraph(line, style="No Spacing")
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif re.match(r"^\d+\.\s+", line):
            doc.add_paragraph(re.sub(r"^\d+\.\s+", "", line), style="List Number")
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(path)
    return path


def maybe_create_pptx(task_id: str, markdown_text: str) -> bool:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return False

    prs = Presentation()
    sections = re.split(r"(?m)^##\s+", markdown_text or "")
    for idx, section in enumerate(sections):
        section = section.strip()
        if not section:
            continue
        lines = section.splitlines()
        title = lines[0].strip() if idx > 0 else "任务结果"
        body = "\n".join(lines[1:] if idx > 0 else lines)[:1200]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    prs.save(os.path.join(artifact_dir, "report.pptx"))
    return True


def save_text_artifacts(task_id: str, answer: str, plan: dict):
    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    os.makedirs(artifact_dir, exist_ok=True)

    if any(x in plan["wanted_artifacts"] for x in ("markdown", "pdf")):
        report_md = os.path.join(artifact_dir, "report.md")
        if not os.path.exists(report_md):
            with open(report_md, "w", encoding="utf-8") as f:
                f.write((answer or "").strip() + "\n")

    if "csv" in plan["wanted_artifacts"]:
        csv_text = extract_fenced_block(answer, "csv")
        if csv_text:
            rows = list(csv.reader(csv_text.splitlines()))
            save_csv_rows(os.path.join(artifact_dir, "result.csv"), rows)
        else:
            table_text = extract_markdown_table(answer)
            if table_text:
                rows = markdown_table_to_rows(table_text)
                if rows:
                    save_csv_rows(os.path.join(artifact_dir, "result.csv"), rows)

    if "json" in plan["wanted_artifacts"]:
        json_text = extract_fenced_block(answer, "json")
        if json_text:
            with open(os.path.join(artifact_dir, "result.json"), "w", encoding="utf-8") as f:
                f.write(json_text + "\n")

    if "xlsx" in plan["wanted_artifacts"]:
        rows = []
        csv_path = os.path.join(artifact_dir, "result.csv")
        if os.path.exists(csv_path):
            with open(csv_path, encoding="utf-8", newline="") as f:
                rows = list(csv.reader(f))
        else:
            csv_text = extract_fenced_block(answer, "csv")
            if csv_text:
                rows = list(csv.reader(csv_text.splitlines()))
            else:
                table_text = extract_markdown_table(answer)
                if table_text:
                    rows = markdown_table_to_rows(table_text)
        if rows:
            maybe_create_xlsx_from_rows(os.path.join(artifact_dir, "result.xlsx"), rows)

    if "docx" in plan["wanted_artifacts"]:
        maybe_create_docx(task_id, answer)

    if "pptx" in plan["wanted_artifacts"]:
        maybe_create_pptx(task_id, answer)


def parse_hermes_response(raw_text: str) -> str:
    text = (raw_text or "").strip()
    if not text:
        return ""
    try:
        data = json.loads(text)
    except Exception:
        return text

    if isinstance(data, dict):
        for key in ("result", "output", "text", "message", "content", "stdout"):
            value = data.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
        choices = data.get("choices")
        if isinstance(choices, list) and choices:
            first = choices[0]
            if isinstance(first, dict):
                msg = first.get("message")
                if isinstance(msg, dict) and isinstance(msg.get("content"), str):
                    return msg["content"].strip()
    return text


def call_hermes_bridge(task: dict, execution_prompt: str, timeout: int = 300) -> str:
    payload = {"task": execution_prompt}
    attachment_path = task.get("attachment_path")
    if attachment_path and os.path.exists(attachment_path):
        payload["attachments"] = [{
            "path": attachment_path,
            "name": task.get("attachment_name") or os.path.basename(attachment_path),
            "kind": task.get("attachment_kind") or "file",
        }]

    body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        HERMES_BRIDGE_RUN_URL,
        data=body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        parsed = parse_hermes_response(raw)
        return parsed or "[错误] Hermes 未返回有效内容"
    except urllib.error.HTTPError as e:
        detail = e.read().decode("utf-8", errors="replace") if hasattr(e, "read") else str(e)
        return f"[错误] Hermes bridge HTTP {e.code}: {detail[:300]}"
    except Exception as e:
        return f"[错误] Hermes bridge 调用失败: {e}"


def build_execution_prompt(task_id: str, prompt: str, plan: dict) -> str:
    artifact_dir = os.path.join(ARTIFACTS_ROOT, task_id)
    wanted = "、".join(plan["wanted_artifacts"])
    base = [
        "你在处理 openclaw-web 的后台任务模式请求。",
        f"任务目标：{prompt}",
        f"用户期望的附件类型：{wanted}",
        "先判断用户真正需要的交付形式，不要默认一律长报告。",
    ]

    if "markdown" in plan["wanted_artifacts"] or "pdf" in plan["wanted_artifacts"]:
        base.append(
            f"如果适合文档交付，请输出完整 Markdown 正文，并确保内容可保存到 {artifact_dir}/report.md。"
        )
    if "csv" in plan["wanted_artifacts"]:
        base.append("如果适合表格交付，请额外提供 ```csv fenced code block```。")
    if "json" in plan["wanted_artifacts"]:
        base.append("如果适合结构化交付，请额外提供 ```json fenced code block```。")
    if "xlsx" in plan["wanted_artifacts"]:
        base.append("如果用户想要 Excel，请至少给出完整的 ```csv fenced code block```，便于后续转换为 xlsx。")
    if "docx" in plan["wanted_artifacts"]:
        base.append("正文请使用清晰的 Markdown 标题和段落结构，便于转换为 Word 文档。")
    if "pptx" in plan["wanted_artifacts"]:
        base.append("如果用户想要演示稿，请按“二级标题 + 简洁要点”组织内容，便于生成幻灯片。")
    if plan["is_report"]:
        base.append("如果用户是在要深度研究/行业分析，再扩成正式报告；否则优先给正确格式而不是凑字数。")

    return "\n".join(base)


def run_openclaw_agent(session_id: str, prompt: str, timeout: int = 300) -> str:
    try:
        result = subprocess.run(
            OPENCLAW_CMD + ["--session-id", session_id, "--message", prompt],
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        return result.stdout if result.returncode == 0 else f"[错误] {result.stderr[:200]}"
    except Exception as e:
        return f"[错误] {e}"


def process_one():
    conn = get_db(DB_PATH)
    row = conn.execute(
        "SELECT * FROM tasks WHERE status='pending' ORDER BY created_at ASC LIMIT 1"
    ).fetchone()
    if not row:
        conn.close()
        return

    task = dict(row)
    task_id = task["id"]
    delivery = json.loads(task.get("delivery_target") or "{}")
    conv_id = delivery.get("conversation_id")
    print(f"[worker] 处理 {task_id}")

    raw_prompt = task.get("prompt_text") or ""
    plan = build_task_plan(task)
    execution_prompt = build_execution_prompt(task_id, raw_prompt, plan)
    is_image_request = plan["executor"] == "local_image"

    if is_image_request:
        img_prompt = re.sub(r"^(.*?)(生成|画|给我|帮我).*?$", r"\1", raw_prompt).strip()
        if len(img_prompt) < 5:
            img_prompt = raw_prompt
        try:
            env = os.environ.copy()
            if os.environ.get("VIVIAI_API_KEY"):
                env["VIVIAI_API_KEY"] = os.environ["VIVIAI_API_KEY"]
            result = subprocess.run(
                ["/root/bin/nanobanana_generate.py", img_prompt, "2K", "16:9"],
                capture_output=True,
                text=True,
                timeout=300,
                env=env,
            )
            if result.returncode == 0 and result.stdout.strip():
                try:
                    img_data = json.loads(result.stdout.strip())
                    if img_data.get("ok") and img_data.get("output_path"):
                        img_path = img_data["output_path"]
                        img_id = task_id[:16]
                        static_dir = IMAGES_DIR
                        static_dir.mkdir(exist_ok=True)
                        ext = Path(img_path).suffix or ".jpg"
                        dest = static_dir / f"{img_id}{ext}"
                        import shutil
                        shutil.copy2(img_path, dest)
                        answer = f"[图片已生成](/images/{img_id}{ext})"
                    else:
                        answer = "[图片生成失败: API返回格式错误]"
                except Exception:
                    answer = "[图片生成失败: JSON解析错误]"
            else:
                answer = f"[图片生成失败: {result.stderr[:100] if result.stderr else 'unknown error'}]"
        except Exception as e:
            answer = f"[图片生成异常: {e}]"
    else:
        if plan["executor"] == "hermes":
            answer = call_hermes_bridge(task, execution_prompt, timeout=300)
            if answer.startswith("[错误]"):
                print("[worker] Hermes bridge 失败，回退到本地 openclaw")
                answer = run_openclaw_agent("task-" + task_id, execution_prompt, timeout=300)
        else:
            answer = run_openclaw_agent("task-" + task_id, execution_prompt, timeout=300)

    print(f"[worker] 首轮输出字数: {len(answer)}")

    report_content = load_report_markdown(task_id)
    report_len = len(report_content) if report_content else 0
    if not is_image_request and plan["is_report"] and len(answer) < 3000 and report_len < 3000:
        print("[worker] 字数不足3000，自动补充...")
        supplement = (
            "上述报告字数不足3000字，请继续补充内容，"
            "在每个章节下增加更详细的数据分析、案例说明和趋势预测，"
            "使总字数达到3000字以上。然后将补充后的完整报告（包含所有章节）重新保存到"
            f"{ARTIFACTS_ROOT}/{task_id}/report.md，覆盖原有文件。"
        )
        if plan["executor"] == "hermes":
            answer2 = call_hermes_bridge(task, supplement, timeout=180)
            if answer2.startswith("[错误]"):
                answer2 = run_openclaw_agent("task-" + task_id + "-sup", supplement, timeout=180)
        else:
            answer2 = run_openclaw_agent("task-" + task_id + "-sup", supplement, timeout=180)
        if len(answer2) > len(answer):
            answer = answer2

    answer_clean = re.sub(r"附：LaTeX 源码.*", "", answer, flags=re.DOTALL).rstrip()

    if not is_image_request and not answer_clean.startswith("[错误]"):
        save_text_artifacts(task_id, answer_clean, plan)

    report_content = load_report_markdown(task_id)
    if report_content and len(report_content) >= len(answer_clean):
        answer_clean = report_content

    report_md = os.path.join(ARTIFACTS_ROOT, task_id, "report.md")
    if os.path.exists(report_md) and "pdf" in plan["wanted_artifacts"]:
        pdf_url = markdown_to_pdf(report_md, task_id)
        if pdf_url:
            answer_clean += f"\n\n📕 [下载 PDF 版报告]({pdf_url})"

    for link in find_artifact_links(task_id):
        fname = os.path.basename(link)
        if link not in answer_clean and fname not in answer_clean:
            answer_clean += f"\n\n📄 [下载 {fname}]({link})"

    now = datetime.now(timezone.utc).isoformat()
    conn.execute(
        "UPDATE tasks SET status='done',result_text=?,finished_at=? WHERE id=?",
        (answer_clean, now, task_id),
    )

    if conv_id:
        app_conn = get_db(APP_DB_PATH)
        task_conn = get_db(DB_PATH)
        link_row = task_conn.execute(
            "SELECT assistant_message_id FROM task_message_links WHERE task_id=? AND conversation_id=?",
            (task_id, conv_id),
        ).fetchone()
        task_conn.close()

        if link_row and link_row["assistant_message_id"]:
            app_conn.execute(
                "UPDATE messages SET content=?,status='done' WHERE id=?",
                (answer_clean, link_row["assistant_message_id"]),
            )
        else:
            import uuid

            assistant_msg_id = "msg_" + uuid.uuid4().hex[:24]
            app_conn.execute(
                "INSERT INTO messages(id,conversation_id,role,content,status,created_at) VALUES(?,?,?,?,?,?)",
                (assistant_msg_id, conv_id, "assistant", answer_clean, "done", now),
            )

        app_conn.execute(
            "UPDATE conversations SET status='active',updated_at=? WHERE id=?",
            (now, conv_id),
        )
        app_conn.commit()
        app_conn.close()

    conn.commit()
    conn.close()
    print(f"[worker] {task_id} 完成（字数：{len(answer_clean)}）")


if __name__ == "__main__":
    print("[worker] 启动")
    while True:
        try:
            process_one()
        except Exception as e:
            print(f"[worker] 异常: {e}")
            import traceback
            traceback.print_exc()
        time.sleep(2)
